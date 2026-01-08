import streamlit as st
import validators
import requests
import sqlite3
import hashlib
import random
import os
import smtplib

from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from dotenv import load_dotenv

from fpdf import FPDF
from pptx import Presentation

from langchain_core.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_groq import ChatGroq
from langchain_community.document_loaders import UnstructuredURLLoader

from datetime import datetime, timedelta

# ==================================================
# ENV
# ==================================================
load_dotenv()
EMAIL = os.getenv("EMAIL_ADDRESS")
PASS = os.getenv("EMAIL_PASSWORD")

# ==================================================
# DATABASE
# ==================================================
conn = sqlite3.connect("app.db", check_same_thread=False)
cur = conn.cursor()

cur.execute("""
CREATE TABLE IF NOT EXISTS users (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    username TEXT UNIQUE,
    password TEXT
)
""")

cur.execute("""
CREATE TABLE IF NOT EXISTS history (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    user_id INTEGER,
    summary TEXT,
    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
)
""")
conn.commit()

# ==================================================
# AUTH HELPERS
# ==================================================
def hash_password(p):
    return hashlib.sha256(p.encode()).hexdigest()

def register_user(u, p):
    try:
        cur.execute(
            "INSERT INTO users (username,password) VALUES (?,?)",
            (u, hash_password(p))
        )
        conn.commit()
        return True
    except:
        return False

def login_user(u, p):
    cur.execute(
        "SELECT id FROM users WHERE username=? AND password=?",
        (u, hash_password(p))
    )
    return cur.fetchone()

def save_history(uid, summary):
    cur.execute(
        "INSERT INTO history (user_id,summary) VALUES (?,?)",
        (uid, summary)
    )
    conn.commit()

def get_history(uid):
    cur.execute(
        "SELECT summary, created_at FROM history WHERE user_id=? ORDER BY created_at DESC",
        (uid,)
    )
    return cur.fetchall()

# ==================================================
# EMAIL OTP
# ==================================================
def send_otp_email(to_email, otp):
    msg = MIMEMultipart()
    msg["From"] = EMAIL
    msg["To"] = to_email
    msg["Subject"] = "🔐 AI Study Assistant - Email Verification"

    html = f"""
    <html>
    <body style="background:#f4f6fb;font-family:Arial;padding:30px">
    <div style="max-width:520px;background:white;padding:30px;border-radius:14px">
        <h2 style="color:#4f46e5">📚 AI Study Assistant</h2>
        <p>Your OTP for account verification is:</p>
        <div style="
            background:#eef2ff;
            font-size:32px;
            font-weight:bold;
            text-align:center;
            padding:18px;
            border-radius:12px;
            letter-spacing:6px;
            margin:20px 0;">
            {otp}
        </div>
        <p>This OTP is valid for a short time.</p>
    </div>
    </body>
    </html>
    """
    msg.attach(MIMEText(html, "html"))

    with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
        server.login(EMAIL, PASS)
        server.send_message(msg)

# ==================================================
# SESSION STATE
# ==================================================
if "user" not in st.session_state:
    st.session_state.user = None
if "otp" not in st.session_state:
    st.session_state.otp = None
if "pending_user" not in st.session_state:
    st.session_state.pending_user = None

# ==================================================
# PAGE CONFIG
# ==================================================
st.set_page_config(
    page_title="AI Study Assistant",
    page_icon="🦜",
    layout="wide"
)

# ==================================================
# MODERN UI + ANIMATIONS (NO LOGIC)
# ==================================================
st.markdown("""
<style>
@keyframes fadeIn {
    from {opacity:0; transform:translateY(15px);}
    to {opacity:1; transform:translateY(0);}
}
@keyframes gradientMove {
    0% {background-position:0% 50%;}
    100% {background-position:100% 50%;}
}
.stApp {animation:fadeIn .6s ease;}
.main-title{
    font-size:2.6rem;
    font-weight:800;
    background:linear-gradient(90deg,#6366f1,#22d3ee,#6366f1);
    background-size:200%;
    -webkit-background-clip:text;
    color:transparent;
    animation:gradientMove 4s linear infinite;
}
.card{
    animation:fadeIn .5s ease;
    border-radius:18px;
    padding:24px;
    line-height:1.6;
}
.dark .stApp{
    background:linear-gradient(135deg,#0f172a,#020617);
    color:#e5e7eb;
}
.dark .card{
    background:rgba(255,255,255,0.08);
    border:1px solid rgba(255,255,255,.12);
}
.light .stApp{background:#f5f7fb;}
.light .card{
    background:white;
    box-shadow:0 15px 35px rgba(0,0,0,.08);
}
.stButton>button{
    border-radius:12px;
    font-weight:600;
    transition:.25s;
}
.stButton>button:hover{
    transform:translateY(-2px);
    box-shadow:0 10px 20px rgba(0,0,0,.15);
}
</style>
""", unsafe_allow_html=True)

# ==================================================
# THEME
# ==================================================
theme = st.sidebar.radio("🌗 Theme", ["Dark", "Light"])
st.markdown(f"<div class='{theme.lower()}'>", unsafe_allow_html=True)

# ==================================================
# LOGIN / REGISTER
# ==================================================
if not st.session_state.user:
    st.markdown("<div class='main-title'>🔐 AI Study Assistant</div>", unsafe_allow_html=True)
    t1, t2 = st.tabs(["Login", "Register"])

    with t1:
        u = st.text_input("Username")
        p = st.text_input("Password", type="password")
        if st.button("Login"):
            user = login_user(u, p)
            if user:
                st.session_state.user = {"id": user[0], "username": u}
                st.success("Logged in")
                st.rerun()
            else:
                st.error("Invalid credentials")

    with t2:
        email = st.text_input("📧 Email")
        ru = st.text_input("New Username")
        rp = st.text_input("New Password", type="password")

        if st.button("📨 Send OTP"):
            if not email or not ru or not rp:
                st.error("Fill all fields")
            else:
                otp = random.randint(100000, 999999)
                st.session_state.otp = str(otp)
                st.session_state.pending_user = (ru, rp)
                send_otp_email(email, otp)
                st.success("OTP sent to email")

        otp_input = st.text_input("🔐 Enter OTP")

        if st.button("✅ Verify & Register"):
            if otp_input == st.session_state.otp:
                ru, rp = st.session_state.pending_user
                if register_user(ru, rp):
                    st.success("Account created")
                    st.session_state.otp = None
                    st.session_state.pending_user = None
                else:
                    st.error("Username exists")
            else:
                st.error("Invalid OTP")

    st.stop()

# ==================================================
# LOGOUT
# ==================================================
with st.sidebar:
    st.markdown(f"👤 **{st.session_state.user['username']}**")
    if st.button("🚪 Logout"):
        st.session_state.user = None
        st.rerun()

# ==================================================
# HEADER
# ==================================================
st.markdown("<div class='main-title'>🦜 AI Study Assistant</div>", unsafe_allow_html=True)
st.caption("Login • Email Verified • Summary • MCQs • PDF • PPT")

# ==================================================
# SIDEBAR
# ==================================================
with st.sidebar:
    groq_api_key = st.text_input("🔑 Groq API Key", type="password")
    summary_type = st.selectbox(
        "Summary Mode",
        ["Paragraph Summary", "Bullet Points (Exam Mode)", "Key Points + Conclusion"]
    )
    language = st.selectbox("Language", ["English", "Hindi", "French", "German"])
    gen_mcq = st.checkbox("Generate MCQs")
    gen_kw = st.checkbox("Generate Keywords")
    gen_ppt = st.checkbox("Generate PPT")
    word_limit = st.slider("Word Limit", 100, 600, 300)

url = st.text_input("🔗 Enter Website or YouTube URL")

# ==================================================
# CONTENT LOADER
# ==================================================
def load_content(url):
    if "youtube.com" in url or "youtu.be" in url:
        r = requests.get(
            "https://www.youtube.com/oembed",
            params={"url": url, "format": "json"}
        )
        if r.status_code != 200:
            return ""
        d = r.json()
        return f"Topic: {d['title']}. Explain clearly with examples."
    loader = UnstructuredURLLoader(urls=[url], ssl_verify=False)
    docs = loader.load()
    return " ".join(d.page_content for d in docs)

# ==================================================
# PROMPT
# ==================================================
def summary_prompt():
    base = f"Write in {language}. Limit {word_limit} words."
    if summary_type == "Bullet Points (Exam Mode)":
        return base + " Use exam-oriented bullet points."
    if summary_type == "Key Points + Conclusion":
        return base + " Use key points and a conclusion."
    return base

# ==================================================
# PDF
# ==================================================
def pdf_bytes(text):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=11)
    pdf.multi_cell(0, 8, text)
    return pdf.output(dest="S").encode("latin-1")

# ==================================================
# PPT
# ==================================================
def make_ppt(text):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Summary"
    slide.placeholders[1].text = text[:1500]
    prs.save("summary.pptx")
    return "summary.pptx"

# ==================================================
# GENERATE
# ==================================================
if st.button("🚀 Generate"):
    if not groq_api_key or not validators.url(url):
        st.error("Provide valid URL & API key")
        st.stop()

    llm = ChatGroq(model="llama-3.1-8b-instant", groq_api_key=groq_api_key)
    text = load_content(url)

    with st.spinner("✨ Generating smart content..."):
        chain = (
            PromptTemplate.from_template(summary_prompt() + "\n\n{text}")
            | llm
            | StrOutputParser()
        )
        summary = chain.invoke({"text": text})

    save_history(st.session_state.user["id"], summary)

    t1, t2, t3 = st.tabs(["📘 Summary", "📝 MCQs", "🔑 Keywords"])

    with t1:
        st.markdown(f"<div class='card'>{summary}</div>", unsafe_allow_html=True)
        st.download_button("📄 Download PDF", pdf_bytes(summary), "summary.pdf")
        if gen_ppt:
            with open(make_ppt(summary), "rb") as f:
                st.download_button("📘 Download PPT", f, "summary.pptx")

    with t2:
        if gen_mcq:
            st.markdown(
                f"<div class='card'>{llm.invoke('Generate 5 MCQs from:\\n'+summary).content}</div>",
                unsafe_allow_html=True
            )

    with t3:
        if gen_kw:
            st.markdown(
                f"<div class='card'>{llm.invoke('Extract 10 keywords from:\\n'+summary).content}</div>",
                unsafe_allow_html=True
            )

# ==================================================
# HISTORY
# ==================================================
st.markdown("## 💾 Your History")
for s, d in get_history(st.session_state.user["id"]):
    with st.expander(d):
        st.write(s)

st.markdown("</div>", unsafe_allow_html=True)
